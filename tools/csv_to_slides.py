#!/usr/bin/env python3
"""Convert semicolon-separated CSV rows into a dark-themed PPTX deck.

Usage:
    csv_to_slides input.csv output.pptx
    csv_to_slides input.xlsx output.pptx
"""

from __future__ import annotations

import argparse
import csv
import html
import json
import os
import re
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, List

from bs4 import BeautifulSoup
from PIL import Image
from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x55, 0x55, 0x55)

SCRIPT_DIR = Path(__file__).resolve().parent
REPO_ROOT = SCRIPT_DIR.parent
LOGO_DIR = REPO_ROOT / "assets" / "lance-logos"
LOGO_MAP_PATH = LOGO_DIR / "logo-map.json"


def normalize_lance_key(value: str | None) -> str:
    s = (value or "").strip().upper()
    s = s.replace("&", "AND")
    s = s.replace("‑", "-").replace("–", "-").replace("—", "-")
    s = re.sub(r"[^A-Z0-9]+", "_", s)
    s = re.sub(r"_+", "_", s).strip("_")
    return s


def load_logo_map() -> Dict[str, str]:
    if not LOGO_MAP_PATH.exists():
        return {}
    try:
        data = json.loads(LOGO_MAP_PATH.read_text(encoding="utf-8"))
        if isinstance(data, dict):
            return {str(k): str(v) for k, v in data.items()}
    except Exception:
        pass
    return {}


def resolve_logo_path(lance: str, logo_map: Dict[str, str]) -> Path | None:
    key = normalize_lance_key(lance)
    if not key:
        return None

    candidates = [key]
    if "_AND_" in key:
        candidates.append(key.replace("_AND_", "_"))
    if key.endswith("_AND"):
        candidates.append(key[: -len("_AND")])

    for cand in candidates:
        mapped = logo_map.get(cand)
        if mapped:
            p = (REPO_ROOT / mapped).resolve()
            if p.exists() and p.suffix.lower() in {".png", ".jpg", ".jpeg"}:
                return p

        direct = LOGO_DIR / f"{cand}.png"
        if direct.exists():
            return direct

    return None


def fit_inside(box_w: int, box_h: int, img_w: int, img_h: int):
    if img_w <= 0 or img_h <= 0:
        return 0, 0, box_w, box_h
    scale = min(box_w / img_w, box_h / img_h)
    w = int(img_w * scale)
    h = int(img_h * scale)
    x = int((box_w - w) / 2)
    y = int((box_h - h) / 2)
    return x, y, w, h


def clean_html_text(raw: str | None) -> str:
    if not raw:
        return ""

    raw = html.unescape(str(raw))
    soup = BeautifulSoup(raw, "html.parser")

    # Preserve line breaks before text extraction.
    for br in soup.find_all("br"):
        br.replace_with("\n")

    text = soup.get_text(separator="\n")
    text = re.sub(r"\r\n?", "\n", text)

    # Keep paragraph spacing readable: collapse >2 blank lines to 2.
    lines = [ln.strip() for ln in text.split("\n")]
    out: List[str] = []
    blank_streak = 0
    for ln in lines:
        if not ln:
            blank_streak += 1
            if blank_streak <= 2:
                out.append("")
            continue
        blank_streak = 0
        out.append(ln)

    cleaned = "\n".join(out).strip()
    return cleaned


def parse_date(value: str | None) -> str:
    if value is None:
        return ""
    s = str(value).strip()
    if not s:
        return ""

    # Unix timestamp support (seconds / milliseconds)
    if re.fullmatch(r"\d+(?:\.\d+)?", s):
        num = float(s)
        if num > 1e12:  # likely milliseconds
            num = num / 1000.0
        try:
            dt = datetime.fromtimestamp(num, tz=timezone.utc)
            return dt.strftime("%B %d, %Y")
        except Exception:
            pass

    # ISO-8601 and common date formats
    iso_candidate = s.replace("Z", "+00:00")
    try:
        dt = datetime.fromisoformat(iso_candidate)
        return dt.strftime("%B %d, %Y")
    except Exception:
        pass

    known_formats = [
        "%Y-%m-%d",
        "%Y/%m/%d",
        "%d.%m.%Y",
        "%d/%m/%Y",
        "%m/%d/%Y",
        "%Y-%m-%d %H:%M:%S",
        "%Y/%m/%d %H:%M:%S",
        "%d.%m.%Y %H:%M:%S",
    ]
    for fmt in known_formats:
        try:
            dt = datetime.strptime(s, fmt)
            return dt.strftime("%B %d, %Y")
        except Exception:
            continue

    # Fallback: return original if parsing fails.
    return s


def get_field(row: Dict[str, str], *names: str) -> str:
    # Direct match first.
    for name in names:
        if name in row:
            return (row.get(name) or "").strip()

    # Case-insensitive fallback.
    lowered = {k.lower().strip(): v for k, v in row.items()}
    for name in names:
        key = name.lower().strip()
        if key in lowered:
            return (lowered[key] or "").strip()

    return ""


def add_styled_run(paragraph, text: str, size_pt: int = 12, color: RGBColor = LIGHT_GRAY, bold: bool = False, underline: bool = False):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.underline = underline
    return run


def clamp(value: int, min_value: int, max_value: int) -> int:
    return max(min_value, min(value, max_value))


def trim_to_word_limit(text: str, max_words: int) -> str:
    matches = list(re.finditer(r"[A-Za-zÀ-ÖØ-öø-ÿ0-9']+", text))
    if len(matches) <= max_words:
        return text
    cut_at = matches[max_words - 1].end()
    return text[:cut_at].rstrip(" ,;:-") + "…"


def summarize_for_slide(text: str, max_words: int = 150, target_words: int = 145) -> str:
    """Extractive summary that represents the full text (not just the beginning)."""
    cleaned = re.sub(r"\s+", " ", (text or "").strip())
    if not cleaned:
        return ""

    words = re.findall(r"[A-Za-zÀ-ÖØ-öø-ÿ0-9']+", cleaned)
    if len(words) <= max_words:
        return cleaned

    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", cleaned) if s.strip()]
    if len(sentences) <= 2:
        # Last-resort word-boundary trim only when sentence structure is missing.
        return trim_to_word_limit(cleaned, max_words)

    stopwords = {
        "the", "a", "an", "and", "or", "but", "if", "then", "than", "that", "this", "these", "those",
        "to", "of", "in", "on", "for", "with", "by", "from", "at", "as", "is", "are", "was", "were",
        "be", "been", "being", "it", "its", "their", "them", "they", "he", "she", "his", "her", "we",
        "our", "you", "your", "i", "me", "my", "not", "no", "yes", "can", "could", "should", "would",
    }

    # Word-frequency scoring.
    freq: Dict[str, int] = {}
    for w in re.findall(r"[A-Za-zÀ-ÖØ-öø-ÿ0-9']+", cleaned.lower()):
        if w in stopwords or len(w) <= 2:
            continue
        freq[w] = freq.get(w, 0) + 1

    sentence_scores = []
    for idx, sentence in enumerate(sentences):
        sent_words = re.findall(r"[A-Za-zÀ-ÖØ-öø-ÿ0-9']+", sentence.lower())
        if not sent_words:
            sentence_scores.append((idx, 0.0))
            continue
        score = sum(freq.get(w, 0) for w in sent_words) / max(len(sent_words), 1)
        sentence_scores.append((idx, score))

    # Ensure coverage across the full text: pick one strong sentence from each third.
    selected = set()
    n = len(sentences)
    thirds = [(0, n // 3), (n // 3, (2 * n) // 3), ((2 * n) // 3, n)]
    for start, end in thirds:
        band = [(i, s) for i, s in sentence_scores if start <= i < end]
        if not band:
            continue
        best_idx = max(band, key=lambda x: x[1])[0]
        selected.add(best_idx)

    # Fill remaining slots by score until target/max words.
    ranked = sorted(sentence_scores, key=lambda x: x[1], reverse=True)

    def current_word_count(indexes):
        return sum(len(sentences[i].split()) for i in indexes)

    for idx, _ in ranked:
        if idx in selected:
            continue
        projected = current_word_count(selected | {idx})
        if projected <= max_words:
            selected.add(idx)
        if projected >= target_words:
            break

    # Keep narrative order.
    ordered = sorted(selected)
    summary = " ".join(sentences[i] for i in ordered).strip()

    # Enforce hard max word limit using regex-token counting.
    summary = trim_to_word_limit(summary, max_words)
    return summary


def render_slide(prs: Presentation, row: Dict[str, str], logo_map: Dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Slide background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BLACK

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    margin = Inches(0.5)
    gap = Inches(0.35)

    full_content_w = slide_w - (margin * 2)

    # Header row: 80% title + 20% logo area.
    title_ratio = 0.80
    title_w = int(full_content_w * title_ratio)
    logo_w = full_content_w - title_w

    # Two-column body layout under metadata.
    columns_w = full_content_w - gap
    left_w = int(columns_w * 0.60)
    right_w = columns_w - left_w

    x_left = margin
    x_right = x_left + left_w + gap
    x_logo = x_left + title_w

    header_h = Inches(0.92)
    meta_h = Inches(0.95)
    content_y = margin + header_h + meta_h + Inches(0.12)
    content_h = slide_h - content_y - margin

    project_id = get_field(row, "Id", "ID", "id")
    title = get_field(row, "Title") or "(Untitled)"
    lens = html.unescape(get_field(row, "Associated Lance", "Associated Lens", "Lens")).strip() or "-"
    deliverable = get_field(row, "Associated Deliverable", "Deliverable") or "-"
    publication_date = parse_date(get_field(row, "Publication Date", "Date")) or "-"
    body = clean_html_text(get_field(row, "Text", "Description", "Body"))

    # Header title area (left 80%)
    title_box = slide.shapes.add_textbox(x_left, margin, title_w, header_h)
    tf_title = title_box.text_frame
    tf_title.clear()
    tf_title.word_wrap = True
    tf_title.vertical_anchor = MSO_ANCHOR.TOP
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT
    title_run = p_title.add_run()
    title_run.text = title
    title_font_size = 28 if len(title) <= 70 else 24 if len(title) <= 110 else 21
    title_run.font.size = Pt(title_font_size)
    title_run.font.bold = True
    title_run.font.color.rgb = WHITE

    # Header logo area (right 20%)
    logo_path = resolve_logo_path(lens, logo_map)
    if logo_path:
        try:
            with Image.open(logo_path) as img:
                off_x, off_y, fit_w, fit_h = fit_inside(logo_w, header_h, img.width, img.height)
            slide.shapes.add_picture(str(logo_path), x_logo + off_x, margin + off_y, fit_w, fit_h)
        except Exception:
            pass

    # Metadata box
    meta_box = slide.shapes.add_textbox(x_left, margin + header_h, left_w, meta_h)
    tf_meta = meta_box.text_frame
    tf_meta.clear()
    tf_meta.word_wrap = True

    p1 = tf_meta.paragraphs[0]
    p1.space_after = Pt(2)
    add_styled_run(p1, f"Associated Lance: {lens}")

    p2 = tf_meta.add_paragraph()
    p2.space_after = Pt(2)
    add_styled_run(p2, "Associated Deliverable: ")
    deliverable_run = add_styled_run(p2, deliverable, underline=True)
    if project_id:
        deliverable_run.hyperlink.address = f"https://guilds.reply.com/news/{project_id}"

    p3 = tf_meta.add_paragraph()
    p3.space_after = Pt(0)
    add_styled_run(p3, f"Publication Date: {publication_date}")

    # Body text area
    body_box = slide.shapes.add_textbox(x_left, content_y, left_w, content_h)
    tf_body = body_box.text_frame
    tf_body.clear()
    tf_body.word_wrap = True
    tf_body.margin_left = Pt(2)
    tf_body.margin_right = Pt(2)
    tf_body.margin_top = Pt(2)
    tf_body.margin_bottom = Pt(2)

    # Summarize long content as a full-text extractive summary (not front-truncation).
    fitted_body = summarize_for_slide(body, max_words=145, target_words=140)
    paragraphs = [seg.strip() for seg in re.split(r"\n\s*\n", fitted_body) if seg.strip()]
    if not paragraphs:
        paragraphs = [""]

    for i, text in enumerate(paragraphs[:4]):
        p = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.line_spacing = 1.2
        p.space_after = Pt(5)
        p.alignment = PP_ALIGN.LEFT

    # Right-side placeholder
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_right, content_y, right_w, content_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = DARK_GRAY
    rect.line.color.rgb = MID_GRAY

    tf_rect = rect.text_frame
    tf_rect.clear()
    tf_rect.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_rect = tf_rect.paragraphs[0]
    p_rect.alignment = PP_ALIGN.CENTER
    ph_run = p_rect.add_run()
    ph_run.text = "Image Placeholder"
    ph_run.font.size = Pt(18)
    ph_run.font.bold = True
    ph_run.font.color.rgb = LIGHT_GRAY


def iter_input_rows(input_path: str):
    ext = os.path.splitext(input_path.lower())[1]

    if ext == ".xlsx":
        wb = load_workbook(input_path, read_only=True, data_only=True)
        try:
            ws = wb.active
            rows = ws.iter_rows(values_only=True)

            header_row = None
            for raw_row in rows:
                if raw_row and any(cell is not None and str(cell).strip() for cell in raw_row):
                    header_row = [str(cell).strip() if cell is not None else "" for cell in raw_row]
                    break

            if not header_row:
                raise ValueError("XLSX appears empty or has no header row.")

            for raw_row in rows:
                values = list(raw_row or [])
                if len(values) < len(header_row):
                    values.extend([None] * (len(header_row) - len(values)))

                row = {
                    header_row[i]: ("" if values[i] is None else str(values[i]).strip())
                    for i in range(len(header_row))
                    if header_row[i]
                }
                if not any((v or "").strip() for v in row.values()):
                    continue
                yield row
        finally:
            wb.close()
        return

    # Default path: semicolon-delimited CSV/text export.
    with open(input_path, "r", encoding="utf-8-sig", newline="") as f:
        reader = csv.DictReader(f, delimiter=";")
        if not reader.fieldnames:
            raise ValueError("CSV appears empty or has no header row.")

        for row in reader:
            if not any((v or "").strip() for v in row.values()):
                continue
            yield row


def convert(input_path: str, pptx_path: str) -> int:
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"Input file not found: {input_path}")

    prs = Presentation()
    # Force 16:9 widescreen.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    logo_map = load_logo_map()

    count = 0
    for row in iter_input_rows(input_path):
        render_slide(prs, row, logo_map)
        count += 1

    if count == 0:
        raise ValueError("No data rows found in input file.")

    out_dir = os.path.dirname(os.path.abspath(pptx_path))
    if out_dir and not os.path.exists(out_dir):
        os.makedirs(out_dir, exist_ok=True)

    prs.save(pptx_path)
    return count


def main() -> int:
    parser = argparse.ArgumentParser(
        prog="csv_to_slides",
        description="Convert CSV or XLSX rows into a dark-themed PPTX deck (1 row = 1 slide).",
    )
    parser.add_argument("input_file", help="Path to input .csv/.xls/.xlsx file")
    parser.add_argument("output_pptx", help="Path to output .pptx file")

    args = parser.parse_args()

    try:
        count = convert(args.input_file, args.output_pptx)
    except Exception as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1

    print(f"Created {count} slide(s): {args.output_pptx}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
